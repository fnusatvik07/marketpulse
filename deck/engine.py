"""
cinematic-deck engine — author slides as SVG, render to crisp PNG (rsvg-convert),
assemble full-bleed into a .pptx (python-pptx). Domain-agnostic; copy into a project's
build dir and edit the PALETTE / WORDMARK block. See references/engine-api.md.

Design space 1280x720, rendered at 2x. Requires: rsvg-convert on PATH; Pillow; python-pptx;
fonts (Anton, Barlow, IBM Plex Mono) installed by scripts/setup.sh.
"""
import os, re, subprocess, base64
from functools import lru_cache
from PIL import ImageFont

HERE   = os.path.dirname(os.path.abspath(__file__))
def _find_fonts():
    cands=[os.environ.get("CINEDECK_FONTS"), os.path.join(HERE,"fonts"),
           os.path.expanduser("~/.claude/skills/cinematic-deck/scripts/fonts"),
           os.path.expanduser("~/Library/Fonts"), os.path.expanduser("~/.fonts")]
    for c in cands:
        if c and os.path.exists(os.path.join(c,"Anton-Regular.ttf")): return c
    return os.path.join(HERE,"fonts")   # default; run setup.sh if missing
FONTS  = _find_fonts()                            # PIL measures TTFs from here (must exist)
ICONS  = os.environ.get("CINEDECK_ICONS", os.path.join(HERE, "icons"))  # ensure_icon cache
OUT    = os.environ.get("CINEDECK_OUT", "out")
os.makedirs(OUT, exist_ok=True)
WORDMARK = os.environ.get("CINEDECK_WORDMARK", "")   # footer brand wordmark ("" = none)
W, H, SCALE = 1280, 720, 2

# ---- PALETTE (MarketPulse paper-light brand) --------------------------------------
# Semantics kept from the dark engine, values inverted for a cream/ink/ledger look:
# INK is the page bg (cream); WHITE is the primary TEXT (ink); surfaces get whiter.
INK="#F2EEE3"; NAVY="#FAF8F1"
SURF1="#FBF9F3"; SURF2="#FFFFFF"; SURF3="#FFFFFF"      # elevation: each lighter than last
LINE="#D8D1BF"; BORDER1="#C8C1AC"; BORDER2="#AFa68E"
BLUE="#1C6B46"; BLUED="#145034"; SKYBLUE="#2F8A61"     # primary accent = ledger green
GOLD="#B07C1F"; CREAM="#3E3A2F"; WHITE="#1F1C15"       # WHITE = ink text on paper
MUTE="#8A8474"; CAP="#5C564A"; FAINT="#9A937F"
RED="#A13327"; GREEN="#1C6B46"; LIGHTCHIP="#FFFFFF"

# ---- Fonts -----------------------------------------------------------------------
DECK_FONTS=os.path.join(os.path.dirname(os.path.abspath(__file__)),"fonts")
F_DISPLAY=os.path.join(DECK_FONTS,"Fraunces-Black.ttf")
F_REG =os.path.join(DECK_FONTS,"SchibstedGrotesk-400.ttf");  F_MED=os.path.join(DECK_FONTS,"SchibstedGrotesk-500.ttf")
F_SEMI=os.path.join(DECK_FONTS,"SchibstedGrotesk-600.ttf"); F_BOLD=os.path.join(DECK_FONTS,"SchibstedGrotesk-700.ttf")
FAM_DISPLAY="Fraunces"; FAM_BODY="Schibsted Grotesk"; FAM_MONO="IBM Plex Mono"

# ---- text measurement / wrapping -------------------------------------------------
@lru_cache(maxsize=None)
def _pil(path,size): return ImageFont.truetype(path,int(size))
def tw(s,path,size): return _pil(path,size).getlength(s)
def wrap(text,path,size,maxw):
    out,cur=[],""
    for w_ in str(text).split():
        t=(cur+" "+w_).strip()
        if tw(t,path,size)<=maxw or not cur: cur=t
        else: out.append(cur); cur=w_
    if cur: out.append(cur)
    return out
def esc(s): return str(s).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
def _ff(family,weight):
    if family==FAM_DISPLAY: return F_DISPLAY
    return F_BOLD if weight>=700 else F_SEMI if weight>=600 else F_MED if weight>=500 else F_REG

# ---- primitives ------------------------------------------------------------------
def txt(x,y,s,size,fill=WHITE,family=FAM_BODY,weight=400,anchor="start",spacing=0,opacity=1,filt=None):
    f=f' filter="url(#{filt})"' if filt else ''
    return (f'<text x="{x}" y="{y}" font-family="{family}" font-size="{size}" font-weight="{weight}" '
            f'fill="{fill}" text-anchor="{anchor}" letter-spacing="{spacing}" opacity="{opacity}"{f}>{esc(s)}</text>')
def rich(x,y,segs,anchor="middle"):
    total=sum(tw(s["t"],_ff(s.get("family",FAM_BODY),s.get("weight",400)),s.get("size",24)) for s in segs)
    sx=x-total/2 if anchor=="middle" else x
    spans="".join(f'<tspan fill="{s.get("fill",WHITE)}" font-size="{s.get("size",24)}" '
                  f'font-weight="{s.get("weight",400)}" font-family="{s.get("family",FAM_BODY)}" '
                  f'letter-spacing="{s.get("spacing",0)}">{esc(s["t"])}</tspan>' for s in segs)
    flt=' filter="url(#hglowG)"' if any(s.get("glow") for s in segs) else ''
    return f'<text x="{sx}" y="{y}" text-anchor="start" xml:space="preserve"{flt}>{spans}</text>'
def rrect(x,y,w,h,r,fill="none",stroke="none",sw=1,opacity=1):
    return (f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{r}" ry="{r}" fill="{fill}" '
            f'stroke="{stroke}" stroke-width="{sw}" opacity="{opacity}"/>')
def line(x1,y1,x2,y2,stroke=LINE,sw=1,opacity=1,dash=None):
    d=f' stroke-dasharray="{dash}"' if dash else ""
    return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{stroke}" stroke-width="{sw}" opacity="{opacity}"{d}/>'
def pill(x,y,label,fill="none",stroke=BLUE,text_fill=BLUE,size=15,pad=16,h=30):
    # width must include the 1.5px letter-spacing the text is drawn with
    w=tw(label,F_SEMI,size)+pad*2+1.5*len(label)
    return (rrect(x,y,w,h,h/2,fill=fill,stroke=stroke,sw=1.4)+
            txt(x+w/2,y+h/2+size*0.36,label,size,text_fill,FAM_BODY,600,"middle",1.5)),w
def arrow(x1,y,x2,color=BLUE):
    return line(x1,y,x2-10,y,color,2.5)+f'<path d="M{x2-12} {y-6} L{x2} {y} L{x2-12} {y+6} Z" fill="{color}"/>'
def stat(cx,y,big,label,color=GOLD,size=128):
    return (txt(cx,y,big,size,color,FAM_DISPLAY,400,"middle",filt="hglowG")+
            txt(cx,y+40,label.upper(),20,CAP,FAM_BODY,700,"middle",5))
def eyebrow(x,y,label,color=GOLD):
    return f'<rect x="{x}" y="{y-12}" width="26" height="3" fill="{color}"/>'+txt(x+36,y,label.upper(),16,color,FAM_BODY,700,spacing=3)

@lru_cache(maxsize=None)
def _icon_raw(path):
    with open(path,encoding="utf-8") as f: return f.read()
def icon(key_or_path,x,y,w,h=None):
    path=key_or_path if os.path.exists(key_or_path) else os.path.join(ICONS,key_or_path+".svg")
    raw=_icon_raw(path); m=re.search(r'viewBox="([^"]+)"',raw); vb=m.group(1) if m else "0 0 128 128"
    inner=re.sub(r'</svg>\s*$','',re.sub(r'^.*?<svg[^>]*>','',raw,flags=re.S),flags=re.S)
    h=h or w
    return f'<svg x="{x}" y="{y}" width="{w}" height="{h}" viewBox="{vb}" preserveAspectRatio="xMidYMid meet">{inner}</svg>'
def logo_chip(cx,cy,size,key):
    r=size/2; ins=size*0.74
    return (f'<g filter="url(#ds)"><rect x="{cx-r}" y="{cy-r}" width="{size}" height="{size}" rx="{size*0.22}" fill="{LIGHTCHIP}" stroke="{LINE}" stroke-width="1"/></g>'
            +icon(key,cx-ins/2,cy-ins/2,ins,ins))
def badge(cx,cy,d,key):
    return (f'<circle cx="{cx}" cy="{cy}" r="{d/2}" fill="{LIGHTCHIP}" stroke="{LINE}" stroke-width="1"/>'
            +icon(key,cx-d*0.37,cy-d*0.37,d*0.74,d*0.74))
def ensure_icon(name, coll="logos"):
    """Fetch ANY brand/tech logo by name from Iconify into the local cache; return its key.
    Pick names that match the DECK'S TOPIC — nothing is hardcoded. Examples:
      ensure_icon('stripe'), ensure_icon('react'), ensure_icon('openai-icon'),
      ensure_icon('python'), ensure_icon('tensorflow'), ensure_icon('figma').
    Browse full-colour logo names at https://icon-sets.iconify.design/logos/ .
    Other collections: coll='simple-icons' (monochrome), 'devicon', 'skill-icons'.
    Returns the key for icon()/logo_chip()/badge(), or None if not found."""
    key=f"{coll}__{name}".replace(":","_").replace("/","_")
    path=os.path.join(ICONS,key+".svg")
    if not os.path.exists(path):
        import urllib.request
        urls=[f"https://api.iconify.design/{coll}/{name}.svg",
              f"https://api.simplesvg.com/{coll}/{name}.svg",
              f"https://unpkg.com/@iconify/json/icons/{coll}/{name}.svg"]
        for u in urls:
            try:
                req=urllib.request.Request(u,headers={"User-Agent":"Mozilla/5.0 (cinematic-deck)"})
                data=urllib.request.urlopen(req,timeout=20).read()
                if b"<svg" not in data[:400]: continue
                os.makedirs(ICONS,exist_ok=True); open(path,"wb").write(data); break
            except Exception:
                continue
        else:
            print("ensure_icon: could not fetch",coll+":"+name); return None
    return key
def data_uri(path):
    ext="png" if path.lower().endswith("png") else "jpeg"
    with open(path,"rb") as f: return f"data:image/{ext};base64,"+base64.b64encode(f.read()).decode()

# ---- shared defs / cinematic layer -----------------------------------------------
BASE_DEFS=f'''
<linearGradient id="bg" x1="0" y1="0" x2="0" y2="1"><stop offset="0" stop-color="{NAVY}"/><stop offset="0.55" stop-color="#F6F2E7"/><stop offset="1" stop-color="{INK}"/></linearGradient>
<radialGradient id="sun" cx="0.82" cy="0.86" r="0.55"><stop offset="0" stop-color="{GOLD}" stop-opacity="0.14"/><stop offset="0.45" stop-color="{GOLD}" stop-opacity="0.04"/><stop offset="1" stop-color="{GOLD}" stop-opacity="0"/></radialGradient>
<linearGradient id="blue" x1="0" y1="0" x2="1" y2="1"><stop offset="0" stop-color="{BLUE}"/><stop offset="1" stop-color="{BLUED}"/></linearGradient>
<linearGradient id="cardg" x1="0" y1="0" x2="0" y2="1"><stop offset="0" stop-color="#FFFFFF"/><stop offset="1" stop-color="#FBF8F0"/></linearGradient>
<filter id="soft" x="-50%" y="-50%" width="200%" height="200%"><feGaussianBlur stdDeviation="14"/></filter>
<filter id="ds" x="-30%" y="-30%" width="160%" height="160%"><feDropShadow dx="0" dy="5" stdDeviation="10" flood-color="#3A3422" flood-opacity="0.18"/></filter>
<linearGradient id="scrimL" x1="0" y1="0" x2="1" y2="0"><stop offset="0" stop-color="#15120B" stop-opacity="0.92"/><stop offset="0.45" stop-color="#15120B" stop-opacity="0.66"/><stop offset="1" stop-color="#15120B" stop-opacity="0.1"/></linearGradient>
<linearGradient id="scrimR" x1="1" y1="0" x2="0" y2="0"><stop offset="0" stop-color="#15120B" stop-opacity="0.92"/><stop offset="0.45" stop-color="#15120B" stop-opacity="0.66"/><stop offset="1" stop-color="#15120B" stop-opacity="0.1"/></linearGradient>
<linearGradient id="scrimFull" x1="0" y1="0" x2="0" y2="1"><stop offset="0" stop-color="#15120B" stop-opacity="0.5"/><stop offset="1" stop-color="#15120B" stop-opacity="0.72"/></linearGradient>
<radialGradient id="vig" cx="0.5" cy="0.43" r="0.82"><stop offset="0.55" stop-color="#3A3422" stop-opacity="0"/><stop offset="1" stop-color="#3A3422" stop-opacity="0.14"/></radialGradient>
<filter id="grain"><feTurbulence type="fractalNoise" baseFrequency="0.9" numOctaves="2" stitchTiles="stitch" result="n"/><feColorMatrix in="n" type="matrix" values="0 0 0 0 0  0 0 0 0 0  0 0 0 0 0  0 0 0 0.5 0"/></filter>
<filter id="hglow" x="-50%" y="-50%" width="200%" height="200%"><feDropShadow dx="0" dy="0" stdDeviation="11" flood-color="{BLUE}" flood-opacity="0.25"/></filter>
<filter id="hglowG" x="-50%" y="-50%" width="200%" height="200%"><feDropShadow dx="0" dy="0" stdDeviation="10" flood-color="{GOLD}" flood-opacity="0.28"/></filter>
<filter id="hglowR" x="-50%" y="-50%" width="200%" height="200%"><feDropShadow dx="0" dy="0" stdDeviation="12" flood-color="{RED}" flood-opacity="0.3"/></filter>
<linearGradient id="gradeP" x1="0" y1="0" x2="1" y2="1"><stop offset="0" stop-color="#EFE9DB" stop-opacity="0.55"/><stop offset="0.5" stop-color="#EFE9DB" stop-opacity="0.08"/><stop offset="1" stop-color="{GOLD}" stop-opacity="0.1"/></linearGradient>
'''
def frame_ticks(m=26,col=GOLD,op=0.5,L=20):
    s=[f'<rect x="{m}" y="{m}" width="{W-2*m}" height="{H-2*m}" fill="none" stroke="{LINE}" stroke-width="1" opacity="0.22"/>']
    for cx,cy,dx,dy in [(m,m,1,1),(W-m,m,-1,1),(m,H-m,1,-1),(W-m,H-m,-1,-1)]:
        s.append(f'<path d="M{cx+dx*L} {cy} L{cx} {cy} L{cx} {cy+dy*L}" stroke="{col}" stroke-width="2" fill="none" opacity="{op}"/>')
    return "\n".join(s)
def footer(topic="",page_no=None):
    y=H-30; p=[line(80,H-58,W-80,H-58,LINE,1,0.45)]
    if WORDMARK:
        p+=[f'<circle cx="85" cy="{y-6}" r="5" fill="{GOLD}"/>',txt(99,y,WORDMARK,18,WHITE,FAM_BODY,700,spacing=0.5)]
    if topic: p.append(txt(W-80,y,topic.upper(),14,MUTE,FAM_BODY,600,"end",2.5))
    if page_no is not None: p.append(txt(W/2,y,str(page_no),13,FAINT,FAM_BODY,500,"middle"))
    return "\n".join(p)
def page(body,defs="",bg="url(#bg)",glow=True,topic="",page_no=None,cine=True):
    layers=[f'<rect width="{W}" height="{H}" fill="{INK}"/>']
    if bg: layers.append(f'<rect width="{W}" height="{H}" fill="{bg}"/>')
    if glow: layers.append(f'<rect width="{W}" height="{H}" fill="url(#sun)"/>')
    over=(f'<rect width="{W}" height="{H}" fill="url(#vig)"/><rect width="{W}" height="{H}" filter="url(#grain)" opacity="0.16"/>'+frame_ticks()) if cine else ""
    return (f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">'
            f'<defs>{BASE_DEFS}{defs}</defs>{"".join(layers)}\n{body}\n{over}\n{footer(topic,page_no)}</svg>')
def render(svg,name):
    sp=os.path.join(OUT,name+".svg"); pp=os.path.join(OUT,name+".png")
    open(sp,"w",encoding="utf-8").write(svg)
    subprocess.run(["rsvg-convert","-w",str(W*SCALE),"-h",str(H*SCALE),"-o",pp,sp],check=True)
    return pp

# ---- glyphs (line-art, brand-coloured) -------------------------------------------
def g_server(cx,cy,s=1.0,color=BLUE):
    w,h=96*s,116*s; x,y=cx-w/2,cy-h/2; u=h/3; out=[rrect(x,y,w,h,10*s,fill="#0C2138",stroke=color,sw=2)]
    for i in range(3):
        uy=y+i*u+6*s; out+=[rrect(x+10*s,uy,w-20*s,u-12*s,4*s,fill="#0A1B2E",stroke=LINE,sw=1.2),
            f'<circle cx="{x+22*s}" cy="{uy+(u-12*s)/2}" r="{3*s}" fill="{GREEN}"/>',
            f'<circle cx="{x+34*s}" cy="{uy+(u-12*s)/2}" r="{3*s}" fill="{GOLD}"/>',
            line(x+w-40*s,uy+(u-12*s)/2,x+w-16*s,uy+(u-12*s)/2,MUTE,2)]
    return "\n".join(out)
def g_cloud(cx,cy,s=1.0,color=BLUE):
    return (f'<g transform="translate({cx},{cy}) scale({s})" fill="#0C2138" stroke="{color}" stroke-width="2.4" stroke-linejoin="round">'
            f'<path d="M-46 18 a26 26 0 0 1 6 -50 a30 30 0 0 1 57 -6 a22 22 0 0 1 9 56 Z"/></g>')
def g_box(cx,cy,s=1.0,color=BLUE,fill="#0E2742"):
    w=92*s; x,y=cx-w/2,cy-w/2
    return (rrect(x,y,w,w,8*s,fill=fill,stroke=color,sw=2.2)+line(x,y+w*0.32,x+w,y+w*0.32,color,1.6,0.7)
            +line(x+w*0.5,y,x+w*0.5,y+w*0.32,color,1.6,0.7))
def g_thermostat(cx,cy,s=1.0):
    import math; r=58*s; out=[f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="#0C2138" stroke="{BLUE}" stroke-width="3"/>',
        f'<circle cx="{cx}" cy="{cy}" r="{r-12*s}" fill="none" stroke="{LINE}" stroke-width="2"/>']
    for a in range(-120,121,24):
        rad=math.radians(a-90); out.append(line(cx+(r-16*s)*math.cos(rad),cy+(r-16*s)*math.sin(rad),cx+(r-24*s)*math.cos(rad),cy+(r-24*s)*math.sin(rad),MUTE,2,0.8))
    rad=math.radians(35-90); out+=[line(cx,cy,cx+(r-22*s)*0.96*math.cos(rad),cy+(r-22*s)*0.96*math.sin(rad),GOLD,4),
        f'<circle cx="{cx}" cy="{cy}" r="{6*s}" fill="{GOLD}"/>']
    return "\n".join(out)
def g_bulb(cx,cy,s=1.0,color=GOLD):
    import math; r=32*s; b=[f'<circle cx="{cx}" cy="{cy-6*s}" r="{r}" fill="#0C2138" stroke="{color}" stroke-width="2.6"/>',
        rrect(cx-13*s,cy+r-12*s,26*s,22*s,4*s,fill="#0C2138",stroke=color,sw=2.4),
        line(cx-9*s,cy-6*s,cx-1*s,cy-16*s,color,2.4),line(cx-1*s,cy-16*s,cx+7*s,cy-6*s,color,2.4)]
    for a in range(0,360,45):
        rad=math.radians(a); b.append(line(cx+(r+8*s)*math.cos(rad),cy-6*s+(r+8*s)*math.sin(rad),cx+(r+20*s)*math.cos(rad),cy-6*s+(r+20*s)*math.sin(rad),color,2.4,0.85))
    return "\n".join(b)
def g_lock(cx,cy,s=1.0,color=GOLD):
    w=66*s;h=50*s;x=cx-w/2;y=cy-h/2+12*s;sh=18*s
    return (rrect(x,y,w,h,8*s,fill="#0C2138",stroke=color,sw=2.6)
            +f'<path d="M{cx-sh} {y} v{-12*s} a{sh} {sh} 0 0 1 {2*sh} 0 v{12*s}" fill="none" stroke="{color}" stroke-width="2.6"/>'
            +f'<circle cx="{cx}" cy="{y+h/2}" r="{5*s}" fill="{color}"/>')
def g_counter(cx,cy,s=1.0):
    return (txt(cx-78*s,cy+16*s,"10",int(62*s),GOLD,FAM_DISPLAY,400,"middle",filt="hglowG")+txt(cx-78*s,cy+44*s,"WISH",int(13*s),CAP,FAM_BODY,700,"middle",3)
            +txt(cx,cy+8*s,"≠",int(40*s),WHITE,FAM_DISPLAY,400,"middle")
            +txt(cx+78*s,cy+16*s,"9",int(62*s),RED,FAM_DISPLAY,400,"middle",filt="hglowR")+txt(cx+78*s,cy+44*s,"NOW",int(13*s),CAP,FAM_BODY,700,"middle",3))
def g_warning(cx,cy,s=1.0,color=RED):
    return (f'<g transform="translate({cx},{cy}) scale({s})" fill="none" stroke="{color}" stroke-width="3" stroke-linejoin="round">'
            f'<path d="M0 -30 L34 30 L-34 30 Z" fill="#2a1116"/><line x1="0" y1="-6" x2="0" y2="12" stroke="{color}"/><circle cx="0" cy="22" r="1.6" fill="{color}" stroke="none"/></g>')
def g_building(cx,cy,s=1.0,sunk=False):
    col=SURF2 if not sunk else "#16314d"; out=[]
    for bx,bh,bw in [(-46,150,28),(-12,200,34),(28,130,30)]:
        x=cx+bx*s; y=cy+(40 if sunk else 0)*s-bh*s/2; out.append(rrect(x,y,bw*s,bh*s,2,fill="#0c2236",stroke=col,sw=1.6))
    return "\n".join(out)
def g_crowd(cx,cy,s=1.0,n=7,color=MUTE):
    out=[]; gap=22*s; start=cx-(n-1)*gap/2
    for i in range(n):
        x=start+i*gap; out+=[f'<circle cx="{x}" cy="{cy-10*s}" r="{7*s}" fill="{color}"/>',f'<path d="M{x-9*s} {cy+18*s} a{9*s} {12*s} 0 0 1 {18*s} 0 Z" fill="{color}"/>']
    return "\n".join(out)
GLYPH={"server":g_server,"cloud":g_cloud,"box":g_box,"thermostat":g_thermostat,"bulb":g_bulb,
       "lock":g_lock,"counter":g_counter,"warning":g_warning,"building":g_building,"crowd":g_crowd}

# ---- composite helper ------------------------------------------------------------
def card(x,y,w,h,label,sub="",icon_key=None,glyph=None,accent=BLUE):
    b=[f'<g filter="url(#ds)">'+rrect(x,y,w,h,16,fill="url(#cardg)",stroke=BORDER1,sw=1.6)+'</g>']; cx=x+w/2
    if icon_key: b.append(logo_chip(cx,y+58,80,icon_key))
    elif glyph: b.append(GLYPH[glyph](cx,y+58,0.7,accent))
    b.append(txt(cx,y+h-34,label,21,WHITE,FAM_BODY,700,"middle"))
    if sub: b.append(txt(cx,y+h-12,sub,14,CAP,FAM_BODY,500,"middle"))
    return "\n".join(b)
def content_header(b,kicker,title,ty=128):
    b.append(eyebrow(84,74,kicker))
    for i,ln in enumerate(wrap(title,F_DISPLAY,44,W-300)):
        b.append(txt(84,ty+i*48,ln,44,WHITE,FAM_DISPLAY,400,spacing=0.5))
    return ty+48*len(wrap(title,F_DISPLAY,44,W-300))

# ---- LAYOUTS ---------------------------------------------------------------------
def title_slide(headline_lines,subtitle,kicker,poster=None,topic="",page_no=1):
    b=[f'<rect width="{W}" height="{H}" fill="url(#bg)"/>',f'<rect width="{W}" height="{H}" fill="url(#sun)"/>']
    if poster:
        pw,ph=300,480; px,py=W-pw-70,(H-ph)/2-10
        b+=[rrect(px-6,py-6,pw+12,ph+12,16,fill="#0A1A2E",stroke=LINE,sw=1.5),
            f'<clipPath id="pc"><rect x="{px}" y="{py}" width="{pw}" height="{ph}" rx="12"/></clipPath>',
            f'<image href="{poster}" x="{px}" y="{py}" width="{pw}" height="{ph}" preserveAspectRatio="xMidYMid slice" clip-path="url(#pc)"/>']
    lx=84; p,_=pill(lx,96,kicker,stroke=BLUE,text_fill=BLUE); b.append(p); y=220
    for ln in headline_lines: b.append(txt(lx,y,ln,74,WHITE,FAM_DISPLAY,400,spacing=0.5,filt="hglow")); y+=74
    b+=[f'<rect x="{lx}" y="{y-2}" width="120" height="4" fill="{GOLD}"/>',txt(lx,y+44,subtitle,24,CREAM,FAM_BODY,500)]
    return page("\n".join(b),bg=INK,glow=False,topic=topic,page_no=page_no)
def divider(part_label,title_lines,blurb="",topic="",page_no=None):
    b=[eyebrow(84,250,part_label)]; y=320
    for ln in title_lines: b.append(txt(84,y,ln,66,WHITE,FAM_DISPLAY,400,spacing=0.5)); y+=66
    if blurb:
        for i,ln in enumerate(wrap(blurb,F_REG,22,760)): b.append(txt(84,y+18+i*32,ln,22,CAP,FAM_BODY,400))
    b.append(f'<rect x="84" y="{y+(60 if blurb else 20)}" width="90" height="4" fill="{GOLD}"/>')
    return page("\n".join(b),topic=topic,page_no=page_no)
def statement(big_lines,sub="",accent=GOLD,topic="",page_no=None,glow=True):
    b=[]; total=len(big_lines); y=H/2-(total-1)*36-(20 if sub else 0)
    for ln in big_lines:
        acc=ln.startswith("*"); col=accent if acc else WHITE
        b.append(txt(W/2,y,ln.lstrip("*"),58,col,FAM_DISPLAY,400,"middle",0.5,filt=("hglowG" if acc else "hglow"))); y+=74
    if sub:
        for i,ln in enumerate(wrap(sub,F_REG,23,880)): b.append(txt(W/2,y+12+i*32,ln,23,CAP,FAM_BODY,400,"middle"))
    return page("\n".join(b),topic=topic,page_no=page_no,glow=glow)
def statement_icon(big_lines,sub,glyph_svg,topic="",page_no=None,accent=GOLD):
    b=[glyph_svg]; y=408
    for ln in big_lines:
        acc=ln.startswith("*"); col=accent if acc else WHITE
        b.append(txt(W/2,y,ln.lstrip("*"),54,col,FAM_DISPLAY,400,"middle",0.5,filt=("hglowG" if acc else "hglow"))); y+=64
    if sub:
        for i,ln in enumerate(wrap(sub,F_REG,22,880)): b.append(txt(W/2,y+16+i*30,ln,22,CAP,FAM_BODY,400,"middle"))
    return page("\n".join(b),topic=topic,page_no=page_no)
def diagram_slide(kicker,title,diagram_inner,caption="",topic="",page_no=None):
    b=[]; content_header(b,kicker,title); b.append(diagram_inner)
    if caption:
        for i,ln in enumerate(wrap(caption,F_REG,19,W-180)): b.append(txt(W/2,H-92+i*26,ln,19,CAP,FAM_BODY,500,"middle"))
    return page("\n".join(b),topic=topic,page_no=page_no)
def photo_hero(uri,kicker,big_lines,sub="",topic="",page_no=None,align="left",scrim="L"):
    b=[f'<image href="{uri}" x="0" y="0" width="{W}" height="{H}" preserveAspectRatio="xMidYMid slice"/>',
       f'<rect width="{W}" height="{H}" fill="url(#gradeP)"/>',f'<rect width="{W}" height="{H}" fill="url(#scrim{scrim})"/>']
    if align=="center":
        cx=W/2; b+=[f'<rect x="{cx-14}" y="244" width="28" height="3" fill="{GOLD}"/>',txt(cx,268,kicker.upper(),16,GOLD,FAM_BODY,700,"middle",4)]; y=346
        for ln in big_lines: b.append(txt(cx,y,ln,60,WHITE,FAM_DISPLAY,400,"middle",0.5,filt="hglow")); y+=68
        if sub:
            for i,ln in enumerate(wrap(sub,F_REG,23,920)): b.append(txt(cx,y+16+i*32,ln,23,CREAM,FAM_BODY,500,"middle"))
    else:
        tx=88; b.append(eyebrow(tx,248,kicker)); y=318
        for ln in big_lines: b.append(txt(tx,y,ln,58,WHITE,FAM_DISPLAY,400,"start",0.5,filt="hglow")); y+=64
        b.append(f'<rect x="{tx}" y="{y-4}" width="96" height="4" fill="{GOLD}"/>')
        if sub:
            for i,ln in enumerate(wrap(sub,F_REG,22,600)): b.append(txt(tx,y+38+i*30,ln,22,CREAM,FAM_BODY,500))
    return page("\n".join(b),bg=None,glow=False,topic=topic,page_no=page_no)
def split_photo(uri,kicker,title,body_paras,topic="",page_no=None,photo_side="right"):
    pw=int(W*0.46); px=W-pw if photo_side=="right" else 0
    tx=84 if photo_side=="right" else pw+72; txtw=W-pw-tx-48 if photo_side=="right" else W-tx-64
    b=[f'<image href="{uri}" x="{px}" y="0" width="{pw}" height="{H}" preserveAspectRatio="xMidYMid slice"/>',
       f'<rect x="{px}" y="0" width="{pw}" height="{H}" fill="url(#gradeP)"/>']
    b.append(f'<rect x="{px}" y="0" width="180" height="{H}" fill="url(#scrimL)"/>' if photo_side=="right"
             else f'<rect x="{px+pw-180}" y="0" width="180" height="{H}" fill="url(#scrimR)"/>')
    b.append(eyebrow(tx,158,kicker)); yy=216
    for ln in wrap(title,F_DISPLAY,46,txtw): b.append(txt(tx,yy,ln,46,WHITE,FAM_DISPLAY,400,"start",0.5)); yy+=50
    yy+=22
    for para in body_paras:
        for ln in wrap(para,F_REG,21,txtw): b.append(txt(tx,yy,ln,21,CAP,FAM_BODY,400)); yy+=30
        yy+=14
    return page("\n".join(b),bg="url(#bg)",glow=True,topic=topic,page_no=page_no)

# ---- assemble ---------------------------------------------------------------------
def build_pptx(items, out_path, fade=True, jpeg=None):
    """items: list of (name, image_path, is_gif). Each placed full-bleed on a 16:9 slide.
    jpeg=82 -> re-encode still PNGs as JPEG for a much smaller, shareable file (~5-8x smaller);
    GIFs are left untouched so they still animate. Build BOTH (full PNG + jpeg=82) for delivery."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    for _,path,is_gif in items:
        use=path
        if jpeg and not is_gif:
            from PIL import Image
            use=os.path.splitext(path)[0]+".q.jpg"
            Image.open(path).convert("RGB").save(use,"JPEG",quality=int(jpeg),optimize=True)
        s=prs.slides.add_slide(blank); s.shapes.add_picture(use,0,0,width=prs.slide_width,height=prs.slide_height)
        if fade:
            try:
                el=s._element; tr=el.makeelement(qn('p:transition'),{'spd':'med'}); tr.append(el.makeelement(qn('p:fade'),{})); el.append(tr)
            except Exception: pass
    prs.save(out_path); return out_path
def contact_sheet(png_paths, out_path, cols=4, tw_=360, th_=202, pad=16, bg=(8,20,38)):
    from PIL import Image
    rows=(len(png_paths)+cols-1)//cols
    sheet=Image.new("RGB",(cols*tw_+(cols+1)*pad,rows*th_+(rows+1)*pad),bg)
    for i,p in enumerate(png_paths):
        t=Image.open(p).convert("RGB").resize((tw_,th_)); r,c=divmod(i,cols)
        sheet.paste(t,(pad+c*(tw_+pad),pad+r*(th_+pad)))
    sheet.save(out_path); return out_path
